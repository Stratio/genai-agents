"""Tests for pptx_layout module."""

import sys
from pathlib import Path
from unittest.mock import patch, MagicMock

import pytest

# Ensure tools/ is in path
sys.path.insert(0, str(Path(__file__).resolve().parent))

import pptx_layout as pl


# ---------------------------------------------------------------------------
# content_area()
# ---------------------------------------------------------------------------

class TestContentArea:
    def test_without_footer(self):
        area = pl.content_area(has_footer=False)
        assert area["top"] == pl.CONTENT_TOP
        assert area["bottom"] == pl.CONTENT_BOTTOM
        assert area["height"] > 0
        assert area["footer_top"] is None

    def test_with_footer(self):
        area = pl.content_area(has_footer=True)
        assert area["top"] == pl.CONTENT_TOP
        assert area["bottom"] < pl.CONTENT_BOTTOM
        assert area["height"] > 0
        assert area["footer_top"] is not None
        # Footer must be above CONTENT_BOTTOM
        assert area["footer_top"] < pl.CONTENT_BOTTOM

    def test_with_footer_height_positive(self):
        area = pl.content_area(has_footer=True, footer_height=0.5)
        assert area["height"] > 0

    def test_footer_reduces_content_height(self):
        area_no_footer = pl.content_area(has_footer=False)
        area_footer = pl.content_area(has_footer=True)
        assert area_footer["height"] < area_no_footer["height"]


# ---------------------------------------------------------------------------
# chart_area()
# ---------------------------------------------------------------------------

class TestChartArea:
    def test_returns_valid_tuple(self):
        left, top, width, height = pl.chart_area()
        assert left >= pl.MARGIN_LEFT
        assert top >= pl.CONTENT_TOP
        assert width > 0
        assert height > 0

    def test_within_bounds(self):
        left, top, width, height = pl.chart_area()
        assert top + height <= pl.CONTENT_BOTTOM
        assert left + width <= pl.SLIDE_WIDTH - pl.MARGIN_RIGHT + 0.01

    def test_with_footer_smaller(self):
        _, _, _, h_no = pl.chart_area(has_footer=False)
        _, _, _, h_yes = pl.chart_area(has_footer=True)
        assert h_yes < h_no


# ---------------------------------------------------------------------------
# check_bounds()
# ---------------------------------------------------------------------------

class TestCheckBounds:
    def test_valid_element(self):
        assert pl.check_bounds(pl.CONTENT_TOP, 2.0, "test") is True

    def test_overflow_bottom(self, capsys):
        """Reproduce Slide 7 bug: top=6.9, height=0.72 → 7.62 > 7.3."""
        result = pl.check_bounds(6.9, 0.72, "table+note")
        assert result is False
        captured = capsys.readouterr()
        assert "WARNING" in captured.err
        assert "table+note" in captured.err

    def test_above_content_top(self, capsys):
        result = pl.check_bounds(1.0, 0.5, "overlapping-header")
        assert result is False
        captured = capsys.readouterr()
        assert "WARNING" in captured.err

    def test_exactly_at_boundary(self):
        """Element that exactly fills the safe area should pass."""
        height = pl.CONTENT_BOTTOM - pl.CONTENT_TOP
        assert pl.check_bounds(pl.CONTENT_TOP, height, "full") is True


# ---------------------------------------------------------------------------
# fit_content()
# ---------------------------------------------------------------------------

class TestFitContent:
    def test_fits_when_space_available(self):
        result = pl.fit_content(pl.CONTENT_TOP, 2.0)
        assert result == 2.0

    def test_clamps_when_overflow(self):
        """If desired height exceeds available space, clamp it."""
        result = pl.fit_content(6.0, 3.0)
        expected_max = pl.CONTENT_BOTTOM - 6.0
        assert result == pytest.approx(expected_max)

    def test_with_footer_reduces_space(self):
        h_no = pl.fit_content(pl.CONTENT_TOP, 10.0, has_footer=False)
        h_yes = pl.fit_content(pl.CONTENT_TOP, 10.0, has_footer=True)
        assert h_yes < h_no

    def test_zero_when_below_bottom(self):
        result = pl.fit_content(pl.CONTENT_BOTTOM + 1.0, 1.0)
        assert result == 0.0

    def test_slide7_scenario(self):
        """Slide 7: table at 6.9" wants 0.72" → should be clamped to 0.4"."""
        result = pl.fit_content(6.9, 0.72)
        assert result == pytest.approx(pl.CONTENT_BOTTOM - 6.9)
        assert result < 0.72  # was clamped


# ---------------------------------------------------------------------------
# create_presentation()
# ---------------------------------------------------------------------------

class TestCreatePresentation:
    def test_dimensions(self):
        prs, palette = pl.create_presentation("corporate")
        from pptx.util import Inches
        assert prs.slide_width == Inches(10.0)
        assert prs.slide_height == Inches(7.5)

    def test_returns_palette(self):
        _, palette = pl.create_presentation("corporate")
        assert "primary" in palette
        assert "font_main" in palette


# ---------------------------------------------------------------------------
# add_slide_header()
# ---------------------------------------------------------------------------

class TestAddSlideHeader:
    def test_returns_content_top(self):
        prs, palette = pl.create_presentation("corporate")
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        result = pl.add_slide_header(slide, "Test Title", "Subtitle", palette)
        assert result == pl.CONTENT_TOP


# ---------------------------------------------------------------------------
# add_image_safe()
# ---------------------------------------------------------------------------

class TestAddImageSafe:
    def test_image_within_bounds(self, tmp_path):
        """Create a dummy PNG and verify it gets placed within safe area."""
        # Create a minimal 1x1 white PNG
        import struct, zlib
        def _make_png():
            sig = b'\x89PNG\r\n\x1a\n'
            ihdr_data = struct.pack('>IIBBBBB', 1, 1, 8, 2, 0, 0, 0)
            ihdr_crc = zlib.crc32(b'IHDR' + ihdr_data) & 0xffffffff
            ihdr = struct.pack('>I', 13) + b'IHDR' + ihdr_data + struct.pack('>I', ihdr_crc)
            raw = b'\x00\xff\xff\xff'
            compressed = zlib.compress(raw)
            idat_crc = zlib.crc32(b'IDAT' + compressed) & 0xffffffff
            idat = struct.pack('>I', len(compressed)) + b'IDAT' + compressed + struct.pack('>I', idat_crc)
            iend_crc = zlib.crc32(b'IEND') & 0xffffffff
            iend = struct.pack('>I', 0) + b'IEND' + struct.pack('>I', iend_crc)
            return sig + ihdr + idat + iend

        img_path = tmp_path / "test.png"
        img_path.write_bytes(_make_png())

        prs, palette = pl.create_presentation("corporate")
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pl.add_image_safe(slide, str(img_path))
        # If it didn't raise, the image was placed successfully
        assert len(slide.shapes) == 1


# ---------------------------------------------------------------------------
# footer_area()
# ---------------------------------------------------------------------------

class TestAddImageWithAspect:
    """Tests for add_image_with_aspect() aspect-ratio preservation."""

    def _make_png(self, tmp_path, width, height, name="test.png"):
        """Create a PNG of given dimensions using PIL."""
        from PIL import Image
        img = Image.new("RGB", (width, height), color=(255, 255, 255))
        path = tmp_path / name
        img.save(str(path))
        return str(path)

    def test_landscape_limited_by_width(self, tmp_path):
        """Landscape image (1800x700, ratio 2.57) in near-square area (5.6x5.5)
        should be limited by width, with actual_h < max_height."""
        img_path = self._make_png(tmp_path, 1800, 700, "landscape.png")
        prs, _ = pl.create_presentation("corporate")
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        actual_w, actual_h = pl.add_image_with_aspect(
            slide, img_path, left=0.4, top=1.3, max_width=5.6, max_height=5.5
        )
        assert actual_w == pytest.approx(5.6, abs=0.01)
        assert actual_h < 5.5
        # Check aspect ratio preserved
        expected_h = 5.6 / (1800 / 700)
        assert actual_h == pytest.approx(expected_h, abs=0.01)

    def test_portrait_limited_by_height(self, tmp_path):
        """Portrait image (700x1800) in area (5.6x5.5)
        should be limited by height, with actual_w < max_width."""
        img_path = self._make_png(tmp_path, 700, 1800, "portrait.png")
        prs, _ = pl.create_presentation("corporate")
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        actual_w, actual_h = pl.add_image_with_aspect(
            slide, img_path, left=0.4, top=1.3, max_width=5.6, max_height=5.5
        )
        assert actual_h == pytest.approx(5.5, abs=0.01)
        assert actual_w < 5.6
        expected_w = 5.5 * (700 / 1800)
        assert actual_w == pytest.approx(expected_w, abs=0.01)

class TestAddImageSafeAspect:
    """Test that add_image_safe preserves aspect ratio."""

    def test_preserves_aspect_landscape(self, tmp_path):
        """Landscape image should not be distorted by add_image_safe."""
        from PIL import Image
        img = Image.new("RGB", (1800, 700), color=(200, 200, 200))
        img_path = tmp_path / "landscape.png"
        img.save(str(img_path))

        prs, _ = pl.create_presentation("corporate")
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        actual_w, actual_h = pl.add_image_safe(slide, str(img_path))
        # Aspect ratio must be preserved (original ratio = 1800/700 ≈ 2.57)
        original_ratio = 1800 / 700
        actual_ratio = actual_w / actual_h
        assert actual_ratio == pytest.approx(original_ratio, rel=0.01)


class TestFooterArea:
    def test_within_slide(self):
        fa = pl.footer_area()
        assert fa["top"] + fa["height"] <= pl.SLIDE_HEIGHT
        assert fa["top"] >= pl.CONTENT_TOP
