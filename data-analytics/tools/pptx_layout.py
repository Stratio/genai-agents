"""PPTX Layout Helper — Prevents overflow and overlap in PowerPoint slides.

Provides centralized constants, positioning functions, and slide helpers
that ensure content stays within safe areas. Follows the same pattern as
chart_layout.py (matplotlib/Plotly) and css_builder.py (CSS/palette).

Usage:
    from pptx_layout import (
        CONTENT_TOP, CONTENT_BOTTOM, CONTENT_WIDTH,
        content_area, chart_area, footer_area,
        check_bounds, fit_content,
        create_presentation, rgb_color,
        add_slide_header, add_text, add_paragraph,
        fill_shape, add_rect, add_kpi_box,
        add_image_safe, add_image_with_aspect,
        add_footer_note, PANEL_GAP,
    )

    prs, palette = create_presentation("corporate")
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    top = add_slide_header(slide, "Title", "Subtitle", palette)
    add_image_safe(slide, "chart.png")
"""

from __future__ import annotations

import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from css_builder import get_palette

# ---------------------------------------------------------------------------
# A. Constants — slide safe area (all values in inches)
# ---------------------------------------------------------------------------

SLIDE_WIDTH = 10.0
SLIDE_HEIGHT = 7.5

MARGIN_LEFT = 0.4
MARGIN_RIGHT = 0.4
MARGIN_BOTTOM = 0.2

# Header zone: the header (title bar + separator) occupies top to 1.13"
HEADER_ZONE_BOTTOM = 1.13
HEADER_PADDING = 0.17  # clearance between header bottom and content top

# Safe content area
CONTENT_TOP = 1.30      # HEADER_ZONE_BOTTOM + HEADER_PADDING
CONTENT_BOTTOM = 7.30   # SLIDE_HEIGHT - MARGIN_BOTTOM
CONTENT_WIDTH = 9.2     # SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
CONTENT_HEIGHT = 6.0    # CONTENT_BOTTOM - CONTENT_TOP

# Footer defaults
FOOTER_DEFAULT_HEIGHT = 0.35
FOOTER_PADDING = 0.10

# Panel layout
PANEL_GAP = 0.3  # gap estándar entre imagen y panel lateral


# ---------------------------------------------------------------------------
# B. Calculation and validation functions
# ---------------------------------------------------------------------------

def content_area(has_footer: bool = False, footer_height: float = FOOTER_DEFAULT_HEIGHT) -> dict:
    """Return the usable content area as a dict.

    Args:
        has_footer: Whether the slide has a footer note.
        footer_height: Height of the footer zone in inches.

    Returns:
        Dict with keys: top, bottom, left, width, height, footer_top.
    """
    bottom = CONTENT_BOTTOM
    footer_top = None
    if has_footer:
        footer_top = bottom - footer_height - FOOTER_PADDING
        bottom = footer_top - FOOTER_PADDING
    return {
        "top": CONTENT_TOP,
        "bottom": bottom,
        "left": MARGIN_LEFT,
        "width": CONTENT_WIDTH,
        "height": bottom - CONTENT_TOP,
        "footer_top": footer_top,
    }


def chart_area(has_footer: bool = False, footer_height: float = FOOTER_DEFAULT_HEIGHT) -> tuple:
    """Return (left, top, width, height) in inches for placing a chart image.

    The area fits within the safe content zone, accounting for optional footer.
    """
    area = content_area(has_footer, footer_height)
    return (area["left"], area["top"], area["width"], area["height"])


def footer_area() -> dict:
    """Return standard footer position as a dict with top, left, width, height."""
    top = CONTENT_BOTTOM - FOOTER_DEFAULT_HEIGHT
    return {
        "top": top,
        "left": MARGIN_LEFT,
        "width": CONTENT_WIDTH,
        "height": FOOTER_DEFAULT_HEIGHT,
    }


def check_bounds(top: float, height: float, label: str = "element") -> bool:
    """Validate that an element fits within the slide safe area.

    Args:
        top: Top position in inches.
        height: Height in inches.
        label: Descriptive name for warning messages.

    Returns:
        True if within bounds, False if overflow detected.
    """
    bottom = top + height
    ok = True
    if top < CONTENT_TOP:
        print(f"WARNING: {label} starts at {top:.2f}\" which is above CONTENT_TOP ({CONTENT_TOP:.2f}\")",
              file=sys.stderr)
        ok = False
    if bottom > CONTENT_BOTTOM:
        print(f"WARNING: {label} extends to {bottom:.2f}\" which exceeds CONTENT_BOTTOM ({CONTENT_BOTTOM:.2f}\")",
              file=sys.stderr)
        ok = False
    return ok


def fit_content(top: float, desired_height: float, has_footer: bool = False,
                footer_height: float = FOOTER_DEFAULT_HEIGHT) -> float:
    """Return the maximum usable height, clamped to available space.

    Args:
        top: Where the content starts (inches).
        desired_height: How tall the content wants to be.
        has_footer: Whether there is a footer.
        footer_height: Footer height in inches.

    Returns:
        min(desired_height, available_space) — never exceeds bounds.
    """
    area = content_area(has_footer, footer_height)
    available = area["bottom"] - top
    if available <= 0:
        return 0.0
    return min(desired_height, available)


def create_presentation(style: str = "corporate") -> tuple:
    """Create a new Presentation with standard dimensions and return (prs, palette).

    Args:
        style: Visual style name for palette lookup.

    Returns:
        Tuple of (Presentation, palette_dict).
    """
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)
    palette = get_palette(style)
    return prs, palette


def rgb_color(rgb_tuple: tuple) -> RGBColor:
    """Convert an (R, G, B) tuple (0-255) to a pptx RGBColor."""
    return RGBColor(rgb_tuple[0], rgb_tuple[1], rgb_tuple[2])


# ---------------------------------------------------------------------------
# C. Slide helpers
# ---------------------------------------------------------------------------

def add_text(tf, text: str, size: int = 14, bold: bool = False,
             color: tuple | None = None, align: str = "left",
             italic: bool = False) -> None:
    """Format text in a TextFrame's first paragraph.

    Args:
        tf: pptx TextFrame object.
        text: Text content.
        size: Font size in points.
        bold: Whether to bold.
        color: Optional (R, G, B) tuple.
        align: "left", "center", or "right".
        italic: Whether to italicize.
    """
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    if color:
        p.font.color.rgb = rgb_color(color)
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)


def add_paragraph(tf, text: str, size: int = 14, bold: bool = False,
                  color: tuple | None = None, align: str = "left",
                  italic: bool = False) -> None:
    """Add an additional paragraph to an existing TextFrame.

    Args:
        tf: pptx TextFrame object.
        text: Text content.
        size: Font size in points.
        bold: Whether to bold.
        color: Optional (R, G, B) tuple.
        align: "left", "center", or "right".
        italic: Whether to italicize.
    """
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    if color:
        p.font.color.rgb = rgb_color(color)
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)


def fill_shape(shape, color: tuple) -> None:
    """Apply a solid fill to a shape.

    Args:
        shape: pptx shape object.
        color: (R, G, B) tuple.
    """
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb_color(color)


def add_rect(slide, left: float, top: float, w: float, h: float,
             color: tuple) -> object:
    """Add a filled rectangle to a slide.

    Args:
        slide: pptx slide object.
        left, top, w, h: Position and size in inches.
        color: (R, G, B) tuple for fill.

    Returns:
        The created shape.
    """
    from pptx.util import Inches as In
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        In(left), In(top), In(w), In(h)
    )
    fill_shape(shape, color)
    shape.line.fill.background()  # no border
    return shape


def add_kpi_box(slide, left: float, top: float, w: float, h: float,
                label: str, value: str, unit: str = "",
                bg_color: tuple | None = None) -> None:
    """Add a KPI box with a large value and label underneath.

    Args:
        slide: pptx slide object.
        left, top, w, h: Position and size in inches.
        label: KPI name (small text below value).
        value: KPI value (large text).
        unit: Optional unit suffix.
        bg_color: Optional background color (R, G, B).
    """
    from pptx.util import Inches as In

    if bg_color:
        add_rect(slide, left, top, w, h, bg_color)

    # Value text box
    txBox = slide.shapes.add_textbox(In(left), In(top + 0.1), In(w), In(h * 0.55))
    tf = txBox.text_frame
    tf.word_wrap = True
    display = f"{value}{unit}" if unit else value
    add_text(tf, display, size=36, bold=True, align="center")

    # Label text box
    txBox2 = slide.shapes.add_textbox(In(left), In(top + h * 0.55), In(w), In(h * 0.35))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    add_text(tf2, label, size=12, align="center", color=(100, 100, 100))


def add_slide_header(slide, title: str, subtitle: str = "",
                     palette: dict | None = None) -> float:
    """Add a standard header with accent bar and separator to a slide.

    Args:
        slide: pptx slide object.
        title: Main title text.
        subtitle: Optional subtitle text.
        palette: Color palette dict from get_palette().

    Returns:
        CONTENT_TOP — the Y position where content should start.
    """
    from pptx.util import Inches as In

    primary = palette.get("primary", (26, 54, 93)) if palette else (26, 54, 93)
    text_color = palette.get("text", (33, 33, 33)) if palette else (33, 33, 33)
    muted_color = palette.get("text_muted", (100, 100, 100)) if palette else (100, 100, 100)

    # Accent bar at top
    bar = slide.shapes.add_shape(1, In(0), In(0), In(SLIDE_WIDTH), In(0.06))
    fill_shape(bar, primary)
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(In(MARGIN_LEFT), In(0.2), In(CONTENT_WIDTH), In(0.5))
    tf = txBox.text_frame
    add_text(tf, title, size=28, bold=True, color=text_color)

    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(In(MARGIN_LEFT), In(0.7), In(CONTENT_WIDTH), In(0.3))
        tf2 = txBox2.text_frame
        add_text(tf2, subtitle, size=14, color=muted_color)

    # Separator line
    sep = slide.shapes.add_shape(1, In(MARGIN_LEFT), In(HEADER_ZONE_BOTTOM - 0.02),
                                 In(CONTENT_WIDTH), In(0.02))
    fill_shape(sep, primary)
    sep.line.fill.background()

    return CONTENT_TOP


def _image_dimensions(img_path: str) -> tuple:
    """Return (width_px, height_px) of an image file using PIL."""
    from PIL import Image
    with Image.open(img_path) as img:
        return img.size


def add_image_safe(slide, img_path: str, has_footer: bool = False,
                   footer_height: float = FOOTER_DEFAULT_HEIGHT) -> tuple:
    """Place an image within the safe content area, preserving aspect ratio.

    The image is fit within the content area and centered vertically.

    Args:
        slide: pptx slide object.
        img_path: Path to the image file (PNG).
        has_footer: Whether there is a footer note on this slide.
        footer_height: Footer height in inches.

    Returns:
        (actual_width, actual_height) in inches after aspect-ratio fitting.
    """
    left, top, max_width, max_height = chart_area(has_footer, footer_height)
    return add_image_with_aspect(slide, img_path, left, top, max_width, max_height)


def add_image_with_aspect(slide, img_path: str, left: float, top: float,
                          max_width: float, max_height: float) -> tuple:
    """Place an image preserving aspect ratio within max_width x max_height.

    Centers vertically if the image is landscape (wider than tall relative
    to the available area).

    Args:
        slide: pptx slide object.
        img_path: Path to the image file.
        left: Left position in inches.
        top: Top position in inches.
        max_width: Maximum available width in inches.
        max_height: Maximum available height in inches.

    Returns:
        (actual_width, actual_height) in inches for caller to position
        adjacent panels.
    """
    from pptx.util import Inches as In

    img_w, img_h = _image_dimensions(img_path)
    img_aspect = img_w / img_h
    area_aspect = max_width / max_height

    if img_aspect > area_aspect:
        # Limited by width
        actual_w = max_width
        actual_h = max_width / img_aspect
    else:
        # Limited by height
        actual_h = max_height
        actual_w = max_height * img_aspect

    # Center vertically within max_height
    offset_y = (max_height - actual_h) / 2

    slide.shapes.add_picture(
        img_path, In(left), In(top + offset_y), In(actual_w), In(actual_h)
    )
    return (actual_w, actual_h)


def add_footer_note(slide, text: str, palette: dict | None = None) -> None:
    """Add a footer note at the standard footer position.

    Args:
        slide: pptx slide object.
        text: Footer text content.
        palette: Optional palette for color.
    """
    from pptx.util import Inches as In

    fa = footer_area()
    muted = palette.get("text_muted", (130, 130, 130)) if palette else (130, 130, 130)

    txBox = slide.shapes.add_textbox(In(fa["left"]), In(fa["top"]),
                                     In(fa["width"]), In(fa["height"]))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_text(tf, text, size=9, italic=True, color=muted)
