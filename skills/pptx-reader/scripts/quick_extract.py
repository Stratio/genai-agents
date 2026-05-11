#!/usr/bin/env python3
"""Quick extraction from a PowerPoint presentation.

Returns structured Markdown on stdout. Diagnostics go to stderr. Exit code 0
on success, 1 on failure. The extractor chain is: python-pptx -> zipfile
XML walk -> soffice text conversion. The first one that produces non-empty
text wins. Legacy .ppt files are auto-converted via LibreOffice first.

Usage:
    python3 quick_extract.py deck.pptx
    python3 quick_extract.py deck.pptx --no-tables
    python3 quick_extract.py deck.pptx --no-notes
    python3 quick_extract.py deck.pptx --include-hidden
    python3 quick_extract.py deck.pptx --tool python-pptx
    cat deck.pptx | python3 quick_extract.py -
"""
from __future__ import annotations

import argparse
import shutil
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path

SUPPORTED_TOOLS = ("auto", "python-pptx", "unzip", "soffice")

ZW_CHARS = "\u200b\u200c\u200d\u200e\u200f\ufeff\u00ad"
SMART_MAP = {
    "\u2018": "'", "\u2019": "'",
    "\u201c": '"', "\u201d": '"',
    "\u2013": "-", "\u2014": "-",
    "\u2026": "...",
}


def _log(msg: str) -> None:
    print(msg, file=sys.stderr)


def _normalize(text: str) -> str:
    for ch in ZW_CHARS:
        text = text.replace(ch, "")
    for k, v in SMART_MAP.items():
        text = text.replace(k, v)
    return text


def _read_stdin_to_tempfile() -> Path:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    tmp.write(sys.stdin.buffer.read())
    tmp.close()
    return Path(tmp.name)


def _is_pptx_zip(path: Path) -> bool:
    try:
        with zipfile.ZipFile(path) as z:
            return "ppt/presentation.xml" in z.namelist()
    except zipfile.BadZipFile:
        return False


def _convert_legacy_ppt(path: Path) -> Path:
    """Convert a legacy binary .ppt to .pptx via soffice. Returns new path."""
    if shutil.which("soffice") is None:
        raise RuntimeError(
            "File looks like a legacy .ppt but 'soffice' (LibreOffice) is not "
            "installed. Install libreoffice or libreoffice-impress."
        )
    out_dir = Path(tempfile.mkdtemp(prefix="pptx_conv_"))
    subprocess.run(
        ["soffice", "--headless", "--convert-to", "pptx",
         "--outdir", str(out_dir), str(path)],
        check=True, capture_output=True,
    )
    converted = out_dir / (path.stem + ".pptx")
    if not converted.exists():
        raise RuntimeError(f"LibreOffice produced no output for {path}")
    return converted


def _slide_is_hidden(slide) -> bool:
    # python-pptx exposes the slide XML element; hidden slides have show="0"
    element = getattr(slide, "_element", None)
    if element is None:
        element = getattr(slide, "element", None)
    if element is None:
        return False
    return element.get("show") == "0"


def _is_title_placeholder(shape) -> bool:
    try:
        if not shape.is_placeholder:
            return False
    except (AttributeError, ValueError):
        return False
    ph = shape.placeholder_format
    if ph is None or ph.idx != 0:
        # idx 0 is always the title placeholder when present
        if ph is None:
            return False
    # Double-check by type: TITLE, CENTER_TITLE, VERTICAL_TITLE
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER  # type: ignore[import-untyped]
    except ImportError:
        return ph is not None and ph.idx == 0
    return ph.type in (
        PP_PLACEHOLDER.TITLE,
        PP_PLACEHOLDER.CENTER_TITLE,
        PP_PLACEHOLDER.VERTICAL_TITLE,
    )


def _find_title_and_element(slide):
    """Return (title_text, title_element) or (None, None).

    ``title_element`` is the underlying lxml element of the shape that
    contributed the title. Use element identity (``is``) to skip it
    during the shape iteration later — python-pptx creates new shape
    proxy objects on each access, so comparing by proxy instance is
    unreliable.
    """
    title_shape = slide.shapes.title
    if title_shape is not None and title_shape.has_text_frame:
        text = title_shape.text_frame.text.strip()  # type: ignore[attr-defined]
        if text:
            return _normalize(text), title_shape._element
    # Fallback: topmost text frame on the slide
    candidates = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()  # type: ignore[attr-defined]
        if not text:
            continue
        top = shape.top or 0
        candidates.append((top, text, shape._element))
    if not candidates:
        return None, None
    candidates.sort(key=lambda t: t[0])
    return _normalize(candidates[0][1]), candidates[0][2]


def _table_to_markdown(shape) -> str:
    rows = []
    for row in shape.table.rows:
        seen = []
        cells = []
        for cell in row.cells:
            cid = id(cell._tc)
            if cid in seen:
                continue
            seen.append(cid)
            text = _normalize(cell.text.strip()).replace("|", "\\|") or " "
            cells.append(text)
        rows.append(cells)
    if not rows:
        return ""
    max_cols = max(len(r) for r in rows)
    rows = [r + [" "] * (max_cols - len(r)) for r in rows]
    header = "| " + " | ".join(rows[0]) + " |"
    sep = "| " + " | ".join(["---"] * max_cols) + " |"
    body = ["| " + " | ".join(r) + " |" for r in rows[1:]]
    return "\n".join([header, sep, *body])


def _extract_with_python_pptx(
    path: Path, include_tables: bool, include_notes: bool, include_hidden: bool,
) -> str:
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError as exc:
        raise RuntimeError("python-pptx not installed") from exc

    prs = Presentation(str(path))
    parts: list[str] = []
    hidden_count = 0

    cp = prs.core_properties
    title = getattr(cp, "title", None) if cp is not None else None
    if title:
        parts.append(f"# {_normalize(title)}\n")

    for idx, slide in enumerate(prs.slides, start=1):
        if _slide_is_hidden(slide):
            hidden_count += 1
            if not include_hidden:
                continue

        title, title_element = _find_title_and_element(slide)
        heading = f"## Slide {idx}: {title}" if title else f"## Slide {idx}"
        parts.append(heading)

        for shape in slide.shapes:
            if _is_title_placeholder(shape):
                continue
            if title_element is not None and shape._element is title_element:
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:  # type: ignore[attr-defined]
                    text = _normalize(
                        "".join(r.text for r in para.runs)
                    ).strip()
                    if not text:
                        continue
                    level = para.level or 0
                    prefix = "  " * level + "- "
                    parts.append(prefix + text)
            elif shape.has_table and include_tables:
                md = _table_to_markdown(shape)
                if md:
                    parts.append("")
                    parts.append(md)

        if include_notes and slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_raw = notes_frame.text if notes_frame is not None else ""
            notes_text = _normalize(notes_raw or "").strip()
            if notes_text:
                parts.append("")
                for line in notes_text.splitlines():
                    parts.append(f"> Notes: {line}" if line.strip() else ">")

        parts.append("")  # blank line between slides

    if hidden_count:
        _log(f"[python-pptx] {hidden_count} hidden slide(s) "
             f"{'included' if include_hidden else 'skipped'}")

    return "\n".join(parts).strip()


def _extract_with_unzip(
    path: Path, include_tables: bool, include_notes: bool, include_hidden: bool,
) -> str:
    # ``include_tables`` is accepted for signature parity with the
    # python-pptx extractor; the XML walk surfaces all text uniformly,
    # table cells included.
    _ = include_tables
    from lxml import etree  # type: ignore[import-untyped]

    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    parts: list[str] = []
    hidden_count = 0

    with zipfile.ZipFile(path) as z:
        slide_names = sorted(
            n for n in z.namelist()
            if n.startswith("ppt/slides/slide") and n.endswith(".xml")
        )
        for idx, name in enumerate(slide_names, start=1):
            with z.open(name) as f:
                tree = etree.parse(f)
            root = tree.getroot()

            # Detect hidden
            if root.get("show") == "0":
                hidden_count += 1
                if not include_hidden:
                    continue

            # Pull every <a:t> for the slide
            texts = [
                (t.text or "") for t in tree.iterfind(f".//{{{NS_A}}}t")
            ]
            slide_text = _normalize(
                "\n".join(t for t in texts if t.strip())
            ).strip()

            parts.append(f"## Slide {idx}")
            if slide_text:
                for line in slide_text.splitlines():
                    s = line.strip()
                    if s:
                        parts.append(f"- {s}")

            # Notes — read sibling notesSlide*.xml if present
            if include_notes:
                notes_name = f"ppt/notesSlides/notesSlide{idx}.xml"
                if notes_name in z.namelist():
                    with z.open(notes_name) as nf:
                        ntree = etree.parse(nf)
                    note_texts = [
                        (t.text or "") for t in ntree.iterfind(f".//{{{NS_A}}}t")
                    ]
                    notes_text = _normalize(
                        "\n".join(t for t in note_texts if t.strip())
                    ).strip()
                    if notes_text:
                        parts.append("")
                        for line in notes_text.splitlines():
                            parts.append(
                                f"> Notes: {line}" if line.strip() else ">"
                            )

            parts.append("")

    if hidden_count:
        _log(f"[unzip] {hidden_count} hidden slide(s) "
             f"{'included' if include_hidden else 'skipped'}")

    return "\n".join(parts).strip()


def _extract_with_soffice(path: Path, **_ignored) -> str:
    """Last-resort: convert to plain text via LibreOffice."""
    if shutil.which("soffice") is None:
        raise RuntimeError("soffice not installed")
    out_dir = Path(tempfile.mkdtemp(prefix="pptx_txt_"))
    subprocess.run(
        ["soffice", "--headless", "--convert-to", "txt",
         "--outdir", str(out_dir), str(path)],
        check=True, capture_output=True,
    )
    converted = out_dir / (path.stem + ".txt")
    if not converted.exists():
        raise RuntimeError(f"soffice produced no text output for {path}")
    return _normalize(converted.read_text(encoding="utf-8", errors="replace")).strip()


def extract(
    path: Path,
    tool: str = "auto",
    include_tables: bool = True,
    include_notes: bool = True,
    include_hidden: bool = False,
) -> str:
    chain = (
        ["python-pptx", "unzip", "soffice"] if tool == "auto" else [tool]
    )
    errors: list[str] = []
    for name in chain:
        try:
            if name == "python-pptx":
                text = _extract_with_python_pptx(
                    path, include_tables, include_notes, include_hidden,
                )
            elif name == "unzip":
                text = _extract_with_unzip(
                    path, include_tables, include_notes, include_hidden,
                )
            elif name == "soffice":
                text = _extract_with_soffice(path)
            else:
                raise RuntimeError(f"unknown tool '{name}'")
        except Exception as exc:  # noqa: BLE001
            _log(f"[{name}] failed: {exc}")
            errors.append(f"{name}: {exc}")
            continue
        if text.strip():
            _log(f"[{name}] produced {len(text)} chars")
            return text
        _log(f"[{name}] produced empty output, trying next tool")
    raise RuntimeError(
        "All extractors failed or produced empty output. Errors: "
        + "; ".join(errors)
    )


def main() -> int:
    parser = argparse.ArgumentParser(
        description=(__doc__ or "").split("\n\n")[0]
    )
    parser.add_argument(
        "path",
        help="Path to the .pptx file, or '-' to read from stdin.",
    )
    parser.add_argument("--no-tables", action="store_true",
                        help="Skip tables, emit only prose and bullets.")
    parser.add_argument("--no-notes", action="store_true",
                        help="Skip speaker notes.")
    parser.add_argument("--include-hidden", action="store_true",
                        help="Include slides marked as hidden "
                             "(default: skip them).")
    parser.add_argument("--tool", choices=SUPPORTED_TOOLS, default="auto",
                        help="Force a specific extractor (default: auto fallback).")
    args = parser.parse_args()

    if args.path == "-":
        path = _read_stdin_to_tempfile()
    else:
        path = Path(args.path)
        if not path.exists():
            _log(f"error: file not found: {path}")
            return 1

    # Detect legacy .ppt and convert up front
    try:
        if not _is_pptx_zip(path):
            _log("input is not a .pptx ZIP; attempting legacy .ppt conversion")
            path = _convert_legacy_ppt(path)
    except Exception as exc:  # noqa: BLE001
        _log(f"error: cannot read input: {exc}")
        return 1

    try:
        text = extract(
            path,
            tool=args.tool,
            include_tables=not args.no_tables,
            include_notes=not args.no_notes,
            include_hidden=args.include_hidden,
        )
    except Exception as exc:  # noqa: BLE001
        _log(f"error: {exc}")
        return 1

    sys.stdout.write(text)
    sys.stdout.write("\n")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
